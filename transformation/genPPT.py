'''
simple script for making your own ppt through typing
'''
from pptx import Presentation
def handle_input_data(data):
    list=[]
    list = data.split('.')
    out_data=''
    for line in list:
        out_data+=(line+'\n')
    return out_data
counter = 0
prs = Presentation()
ppt_name = ''

while True:

    answer = input("-->Do u want to generate new page of ppt?(y/n)<--\n")
    if(counter==0):
        ppt_name = input("please enter ppt name:\n")
    if(answer.lower()=="n"):
        break
    style = int(input("Please choose your page style:\n"
                  "TITLE_SLIDE = 0 \n"
                  "TITLE_AND_CONTENT = 1 \n"
                  "SECTION_HEADER = 2 \n"
                  "TWO_CONTENT = 3  \n"
                  "COMPARISON = 4 \n"
                  "TITLE_ONLY = 5 \n"
                  "BLANK = 6 \n"
                  "CONTENT_WITH_CAPTION = 7 \n"
                  "PICTURE_WITH_CAPTION = 8 \n"
                  "TITLE_AND_VERTICAL_TEXT = 9 \n"
                  "VERTICAL_TITLE_AND_TEXT = 10 \n"))
    counter+=1
    print("page %d ----------------------------------------------------------------------/" % counter)
    slide = prs.slides.add_slide(prs.slide_layouts[style])
    for shape in slide.shapes:
        print("shape style: -----------------------------------------/")
        phf = shape.placeholder_format
        if shape.is_placeholder:
            print('%d, %s' % (phf.idx, phf.type))
        if shape.has_text_frame:
            print("shape name: %s" %shape.name)
            print("has text frame!")
            data = input("enter content here\n")
            shape.text = handle_input_data(data)
        print("-----------------------------------------------------/")





if(counter>0):
    prs.save(ppt_name+'.pptx')


print("Thanks for using!")
